#!/usr/bin/env python3
"""
PPC Companion — Creative Crew
Full creative team that produces professional training materials.

Crew members:
  1. Research Analyst — topic research, data points, case studies
  2. Presentation Writer — slide scripts, talking points, speaker notes
  3. Document Designer — PDF handouts via HTML/CSS + WeasyPrint (or HTML fallback)
  4. Slide Producer — PPTX presentations via python-pptx
  5. Infographic Designer — PNG infographics via Pillow
  6. Spreadsheet Engineer — XLSX workbooks via openpyxl
  7. Creative Director — reviews all outputs for quality and consistency

Usage:
  OPENROUTER_API_KEY=sk-or-... python creative-crew.py "Sponsored Brands" --phase 2 --module 2.1 -o ./output/

Env vars:
  CREW_MODEL       — override LLM (default: openrouter/deepseek/deepseek-v4-flash)
  OPENROUTER_API_KEY — required
"""
import argparse
import json
import os
import re
import sys
import textwrap
from datetime import datetime
from pathlib import Path
from typing import Optional

# ── Config ──────────────────────────────────────────────────────────────────
LLM_MODEL = os.getenv("CREW_MODEL", "openrouter/deepseek/deepseek-v4-flash")
LLM_KEY = os.getenv("OPENROUTER_API_KEY")
LLM_BASE = "https://openrouter.ai/api/v1"

if not LLM_KEY:
    print("FATAL: Set OPENROUTER_API_KEY env var or pass --api-key", file=sys.stderr)
    sys.exit(1)

os.environ["OPENAI_API_KEY"] = LLM_KEY
os.environ["OPENAI_BASE_URL"] = LLM_BASE

from crewai import Agent, Task, Crew, Process, LLM

# ── Brand constants ─────────────────────────────────────────────────────────
BRAND = {
    "name": "Amazon PPC Manager Training Program",
    "tagline": "Student Workbook",
    "version": "2026",
    "primary": "#1a56db",
    "primary_light": "#3b82f6",
    "accent": "#f59e0b",
    "dark": "#111827",
    "bg": "#f8fafc",
    "surface": "#ffffff",
    "text": "#1f2937",
    "text_muted": "#6b7280",
    "border": "#e5e7eb",
    "success": "#10b981",
    "error": "#ef4444",
}

TOPICS = {
    1: {"title": "Foundations", "modules": [
        ("1.1", "How Amazon Works"), ("1.2", "PPC Metrics Fundamentals"), ("1.3", "Search Term Analysis Basics"),
    ]},
    2: {"title": "Amazon Ads Deep Dive", "modules": [
        ("2.1", "Sponsored Brands Campaign Structure"), ("2.2", "Sponsored Display & DSP Overview"), ("2.3", "Keyword Match Types & Negative Targeting"),
    ]},
    3: {"title": "Advanced Strategy", "modules": [
        ("3.1", "Bidding Strategies & Portfolio Management"), ("3.2", "Campaign Architecture & Segmentation"), ("3.3", "Performance Analysis & Optimization Cycles"),
    ]},
    4: {"title": "Capstone", "modules": [("4.1", "Building a Full PPC Strategy")]},
}

# ── Agents ──────────────────────────────────────────────────────────────────
def make_agents(llm: LLM):
    researcher = Agent(
        role="Senior PPC Research Analyst",
        goal="Deliver precise, current Amazon PPC data with real numbers, case studies, and actionable insights",
        backstory=(
            "You are a senior Amazon advertising analyst with 8+ years experience managing "
            "multi-million dollar PPC portfolios. You provide specific data points, real ACoS "
            "benchmarks, actual campaign examples, and industry trends. You never speculate — "
            "every claim is backed by observable data or widely accepted benchmarks."
        ),
        llm=llm, verbose=False,
    )

    writer = Agent(
        role="Presentation Scriptwriter",
        goal="Write compelling slide-by-slide scripts with speaker notes that a trainer can deliver naturally",
        backstory=(
            "You are a corporate presentation writer specialising in e-commerce education. "
            "You structure content in a 'what → why → how → practice' pattern. Every slide "
            "has a single clear message, a supporting visual concept, and detailed speaker "
            "notes. You write for a Filipino VA audience — clear, practical, no jargon without "
            "explanation."
        ),
        llm=llm, verbose=False,
    )

    infographic = Agent(
        role="Infographic Content Strategist",
        goal="Design the content and data layout for a single-page infographic summarising the module",
        backstory=(
            "You are a visual communication expert. You design infographics that pack "
            "maximum information into a single page. You think in terms of visual hierarchy, "
            "colour-coded sections, icon placeholders, data callouts, and comparison layouts. "
            "Your output describes exactly what goes where, with exact labels and numbers."
        ),
        llm=llm, verbose=False,
    )

    spreadsheet = Agent(
        role="Spreadsheet Engineer",
        goal="Design Excel workbooks with calculators, data tables, and formatting specs",
        backstory=(
            "You build professional Excel tools for Amazon sellers: bid calculators, ACoS "
            "trackers, keyword managers, budget planners. You specify exact cell references, "
            "formulas, conditional formatting rules, data validation, and colour schemes. "
            "Your output is a JSON specification that maps directly to openpyxl calls."
        ),
        llm=llm, verbose=False,
    )

    director = Agent(
        role="Creative Director & Quality Reviewer",
        goal="Review all outputs for accuracy, brand consistency, professional quality, and pedagogical soundness",
        backstory=(
            "You are the final quality gate. You check every number, every formula, every "
            "claim for accuracy. You ensure brand consistency (colours, fonts, tone), "
            "pedagogical flow (scaffolding, assessment alignment), and professional polish. "
            "You flag anything that would embarrass the programme."
        ),
        llm=llm, verbose=False,
    )

    return {"researcher": researcher, "writer": writer, "infographic": infographic,
            "spreadsheet": spreadsheet, "director": director}


# ── Tasks ───────────────────────────────────────────────────────────────────
CONTENT_HINT = json.dumps({
    "module_code": "2.1", "module_title": "...",
    "slide_scripts": [
        {"slide_number": 1, "title": "...", "bullets": ["..."], "visual": "description of chart/diagram", "speaker_notes": "...", "duration_seconds": 60}
    ],
    "handout_content": {
        "sections": [{"heading": "...", "content": "...", "key_takeaways": ["..."]}],
        "glossary": [{"term": "...", "definition": "..."}],
        "checklist": ["..."],
        "formula_reference": [{"name": "...", "formula": "...", "example": "..."}]
    },
    "infographic_spec": {
        "title": "...",
        "sections": [{"label": "...", "items": ["..."], "colour": "#hex"}],
        "callout_stats": [{"value": "...", "label": "..."}]
    },
    "spreadsheet_spec": {
        "sheets": [{"name": "...", "columns": ["..."], "sample_rows": [["..."]], "formulas": {"cell": "formula"}, "notes": "..."}]
    }
}, indent=2)

def make_tasks(agents, phase: int, code: str, topic: str):
    hint = CONTENT_HINT

    research = Task(
        description=(
            f"Deep research for Phase {phase}, Module {code}: {topic}.\n"
            "Deliver:\n"
            "1. Key concepts (5-8) with precise definitions\n"
            "2. Real-world examples with actual numbers (ACoS %, ROAS, spend, revenue)\n"
            "3. Common mistakes VAs make with this topic\n"
            "4. Industry benchmarks for 2025-2026\n"
            "5. 2-3 mini case studies (pet brand, electronics, supplements)\n"
            "6. Glossary of 8-10 terms with clear definitions\n"
            "Format as structured JSON."
        ),
        expected_output="Structured research JSON with all 6 sections",
        agent=agents["researcher"],
    )

    write_content = Task(
        description=(
            f"Write complete training content for Phase {phase}, Module {code}: {topic}.\n"
            f"Output valid JSON matching this structure:\n{hint}\n\n"
            "Requirements:\n"
            "- 10-12 slides with speaker notes (each note 3-5 sentences)\n"
            "- Handout with 4 content sections, glossary, checklist, formula reference\n"
            "- Infographic spec with 4 visual sections and 3 callout stats\n"
            "- Spreadsheet spec with 2 sheets (calculator + data tracker)\n"
            "- All content must reference real PPC metrics\n"
            "Output ONLY valid JSON."
        ),
        expected_output="JSON with slide_scripts, handout_content, infographic_spec, spreadsheet_spec",
        agent=agents["writer"],
        context=[research],
    )

    design_infographic = Task(
        description=(
            f"Create a detailed infographic content plan for Phase {phase}, Module {code}: {topic}.\n"
            "Specify exact text, numbers, icons, layout zones, colour coding. "
            "Think like a Canva designer — every element has a position, size, colour. "
            "Output JSON with: title, subtitle, 4-6 zones, each with items, icon suggestions, colours, and placement. "
            "Include 3 callout statistics with exact numbers."
        ),
        expected_output="JSON infographic specification with exact layout details",
        agent=agents["infographic"],
        context=[research, write_content],
    )

    design_spreadsheet = Task(
        description=(
            f"Design 2 Excel sheets for Phase {phase}, Module {code}: {topic}.\n"
            "Sheet 1: Interactive calculator with formulas, input cells, conditional formatting.\n"
            "Sheet 2: Data tracker/template with sample data and headers.\n"
            "Output JSON with: sheet names, column headers, cell formulas, "
            "conditional formatting rules, data validation, colour scheme. "
            "Specify exact cell references (e.g. B2, C5:C20)."
        ),
        expected_output="JSON spreadsheet specification with exact cell references and formulas",
        agent=agents["spreadsheet"],
        context=[research],
    )

    review = Task(
        description=(
            "Review the research, content, infographic spec, and spreadsheet spec.\n"
            "Check for:\n"
            "1. Metric accuracy (ACoS = spend/revenue, not revenue/spend)\n"
            "2. Brand colour consistency\n"
            "3. Pedagogical flow (scaffolding, no knowledge gaps)\n"
            "4. Spreadsheet formula correctness\n"
            "5. Infographic visual balance\n\n"
            "Output a JSON review with: overall_score (1-10), issues_found (list), "
            "approved_content (the final corrected JSON — merge fixes into the original content output).\n"
            "The approved_content MUST be a single JSON object matching the original structure."
        ),
        expected_output="JSON review with approved_content (corrected full content)",
        agent=agents["director"],
        context=[research, write_content, design_infographic, design_spreadsheet],
    )

    return [research, write_content, design_infographic, design_spreadsheet, review]


# ── File Generators ─────────────────────────────────────────────────────────
def extract_json(text: str) -> Optional[dict]:
    text = re.sub(r"^REVIEW_PASSED\s*\n?", "", text, flags=re.MULTILINE)
    match = re.search(r"```(?:json)?\s*\n(.*?)\n```", text, re.DOTALL)
    if match:
        text = match.group(1)
    try:
        return json.loads(text.strip())
    except json.JSONDecodeError:
        return None


def generate_pptx(content: dict, out_path: Path):
    """Generate a professional PPTX presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    primary = RGBColor(0x1a, 0x56, 0xdb)
    accent = RGBColor(0xf5, 0x9e, 0x0b)
    dark = RGBColor(0x11, 0x18, 0x27)
    white = RGBColor(0xff, 0xff, 0xff)
    muted = RGBColor(0x6b, 0x72, 0x80)
    light_bg = RGBColor(0xf0, 0xf4, 0xf8)

    # Title slide
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = primary
    # Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.get("module_title", content.get("topic", "Training Module"))
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = f"Phase {content.get('phase', '?')} · Module {content.get('module_code', '?')}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xbf, 0xdb, 0xfe)
    p2.alignment = PP_ALIGN.CENTER
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.2), Inches(2.3), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # Tagline
    p3 = tf.add_paragraph()
    p3.text = ""
    tag_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(1))
    tft = tag_box.text_frame
    tft.paragraphs[0].text = "Amazon PPC Manager Training Program · Student Workbook · v2026"
    tft.paragraphs[0].font.size = Pt(14)
    tft.paragraphs[0].font.color.rgb = RGBColor(0x93, 0xc5, 0xfd)
    tft.paragraphs[0].alignment = PP_ALIGN.CENTER

    slides = content.get("slide_scripts", [])
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_bg = slide.background.fill
        slide_bg.solid()
        slide_bg.fore_color.rgb = white

        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Pt(6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = primary
        bar.line.fill.background()

        # Slide number
        num_box = slide.shapes.add_textbox(Inches(12.3), Inches(0.3), Inches(0.8), Inches(0.4))
        ntf = num_box.text_frame
        ntf.paragraphs[0].text = str(s.get("slide_number", ""))
        ntf.paragraphs[0].font.size = Pt(11)
        ntf.paragraphs[0].font.color.rgb = muted

        # Title
        t_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(1))
        ttf = t_box.text_frame
        ttf.word_wrap = True
        ttf.paragraphs[0].text = s.get("title", "")
        ttf.paragraphs[0].font.size = Pt(32)
        ttf.paragraphs[0].font.bold = True
        ttf.paragraphs[0].font.color.rgb = dark

        # Bullets
        bullets = s.get("bullets", [])
        if bullets:
            b_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(4.5))
            btf = b_box.text_frame
            btf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = btf.paragraphs[0]
                else:
                    p = btf.add_paragraph()
                p.text = f"  •  {bullet}"
                p.font.size = Pt(18)
                p.font.color.rgb = dark
                p.space_after = Pt(12)

        # Visual placeholder
        visual = s.get("visual", "")
        if visual:
            v_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(2.2), Inches(3.3), Inches(4.5)
            )
            v_box.fill.solid()
            v_box.fill.fore_color.rgb = light_bg
            v_box.line.color.rgb = RGBColor(0xd1, 0xd5, 0xdb)
            vtf = v_box.text_frame
            vtf.word_wrap = True
            vtf.paragraphs[0].text = f"[ {visual[:80]}... ]" if len(visual) > 80 else f"[ {visual} ]"
            vtf.paragraphs[0].font.size = Pt(10)
            vtf.paragraphs[0].font.color.rgb = muted
            vtf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = s.get("speaker_notes", "")

    # Closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = primary
    cbox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    ctf = cbox.text_frame
    ctf.word_wrap = True
    ctf.paragraphs[0].text = "Questions?"
    ctf.paragraphs[0].font.size = Pt(48)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = white
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = ctf.add_paragraph()
    p2.text = "Review your notes and try the exercises"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xbf, 0xdb, 0xfe)
    p2.alignment = PP_ALIGN.CENTER

    prs.save(str(out_path))


def generate_pdf_handout(content: dict, out_path: Path):
    """Generate a professional PDF handout using WeasyPrint. Returns actual path."""
    handout = content.get("handout_content", {})
    sections = handout.get("sections", [])
    glossary = handout.get("glossary", [])
    checklist = handout.get("checklist", [])
    formulas = handout.get("formula_reference", [])

    sections_html = ""
    for sec in sections:
        takeaways = "".join(f"<li>{t}</li>" for t in sec.get("key_takeaways", []))
        sections_html += f"""
        <div class="section">
            <h2>{sec.get('heading', '')}</h2>
            <p>{sec.get('content', '')}</p>
            {f'<div class="takeaway-box"><strong>Key Takeaways</strong><ul>{takeaways}</ul></div>' if takeaways else ''}
        </div>"""

    glossary_html = ""
    if glossary:
        rows = "".join(f"<tr><td><strong>{g.get('term','')}</strong></td><td>{g.get('definition','')}</td></tr>" for g in glossary)
        glossary_html = f"""
        <div class="section">
            <h2>Glossary</h2>
            <table class="glossary">{rows}</table>
        </div>"""

    checklist_html = ""
    if checklist:
        items = "".join(f"<li>{c}</li>" for c in checklist)
        checklist_html = f"""
        <div class="section checklist-box">
            <h2>Quick Reference Checklist</h2>
            <ul class="checklist">{items}</ul>
        </div>"""

    formula_html = ""
    if formulas:
        rows = "".join(f"<tr><td><strong>{f.get('name','')}</strong></td><td class='mono'>{f.get('formula','')}</td><td>{f.get('example','')}</td></tr>" for f in formulas)
        formula_html = f"""
        <div class="section">
            <h2>Formula Reference</h2>
            <table class="formulas">{rows}</table>
        </div>"""

    html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
@page {{ size: A4; margin: 2cm; @bottom-center {{ content: counter(page) " / " counter(pages); font-size: 9px; color: #999; }} }}
body {{ font-family: 'Helvetica Neue', Arial, sans-serif; color: #1f2937; line-height: 1.6; font-size: 11pt; }}
.header {{ border-bottom: 3px solid #1a56db; padding-bottom: 12px; margin-bottom: 24px; }}
.header h1 {{ margin: 0; font-size: 22pt; color: #1a56db; }}
.header .subtitle {{ color: #6b7280; font-size: 10pt; margin-top: 4px; }}
.section {{ margin-bottom: 20px; }}
.section h2 {{ font-size: 14pt; color: #1a56db; border-left: 4px solid #f59e0b; padding-left: 10px; margin-bottom: 8px; }}
.section p {{ margin: 6px 0; }}
.takeaway-box {{ background: #eff6ff; border-radius: 6px; padding: 10px 14px; margin-top: 8px; }}
.takeaway-box ul {{ margin: 4px 0 0 16px; }}
.glossary {{ width: 100%; border-collapse: collapse; font-size: 10pt; }}
.glossary td {{ padding: 6px 8px; border-bottom: 1px solid #e5e7eb; vertical-align: top; }}
.checklist-box {{ background: #f0fdf4; border: 1px solid #bbf7d0; border-radius: 6px; padding: 14px; }}
.checklist {{ margin: 6px 0 0 16px; }}
.formulas {{ width: 100%; border-collapse: collapse; font-size: 10pt; }}
.formulas td {{ padding: 6px 8px; border-bottom: 1px solid #e5e7eb; }}
.mono {{ font-family: 'Courier New', monospace; font-size: 9pt; color: #1a56db; background: #f8fafc; padding: 2px 6px; border-radius: 3px; }}
.footer {{ margin-top: 30px; padding-top: 10px; border-top: 1px solid #e5e7eb; font-size: 8pt; color: #9ca3af; text-align: center; }}
</style>
</head>
<body>
<div class="header">
    <h1>{content.get('module_title', content.get('topic', 'Training Module'))}</h1>
    <div class="subtitle">Phase {content.get('phase', '?')} · Module {content.get('module_code', '?')} · Amazon PPC Manager Training Program · v2026</div>
</div>
{sections_html}
{glossary_html}
{checklist_html}
{formula_html}
<div class="footer">Amazon PPC Manager Training Program · Student Workbook · v2026 · Generated {datetime.now().strftime('%Y-%m-%d')}</div>
</body>
</html>"""

    try:
        from weasyprint import HTML
        HTML(string=html).write_pdf(str(out_path))
    except (ImportError, OSError) as e:
        html_path = out_path.with_suffix('.html')
        html_path.write_text(html)
        print(f"  (PDF engine unavailable: {e} — saved HTML: {html_path})", file=sys.stderr)
        return html_path
    return out_path


def generate_infographic(content: dict, out_path: Path):
    """Generate a PNG infographic using Pillow."""
    from PIL import Image, ImageDraw, ImageFont

    W, H = 2480, 3508  # A4 at 300dpi
    img = Image.new("RGB", (W, H), BRAND["bg"])
    draw = ImageDraw.Draw(img)

    try:
        font_lg = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
        font_md = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 28)
        font_sm = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 22)
        font_xs = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
        font_stat = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 56)
    except OSError:
        font_lg = font_md = font_sm = font_xs = font_stat = ImageFont.load_default()

    colours = ["#1a56db", "#7c3aed", "#059669", "#d97706", "#dc2626"]

    # Header bar
    draw.rectangle([(0, 0), (W, 200)], fill=BRAND["primary"])
    draw.text((80, 50), content.get("module_title", content.get("topic", "")), fill="white", font=font_lg)
    draw.text((80, 130), f"Phase {content.get('phase', '?')} · Module {content.get('module_code', '?')} · Amazon PPC Training v2026",
              fill="#93c5fd", font=font_xs)

    # Callout stats
    stats = content.get("infographic_spec", {}).get("callout_stats", [])
    if stats:
        y = 260
        stat_w = (W - 160) // min(len(stats), 3)
        for i, st in enumerate(stats[:3]):
            x = 80 + i * stat_w
            draw.rounded_rectangle([(x, y), (x + stat_w - 30, y + 220)], radius=16, fill="white", outline=BRAND["border"])
            val = st.get("value", "")
            lbl = st.get("label", "")
            bbox = draw.textbbox((0, 0), val, font=font_stat)
            tw = bbox[2] - bbox[0]
            draw.text((x + (stat_w - 30 - tw) // 2, y + 20), val, fill=BRAND["primary"], font=font_stat)
            bbox2 = draw.textbbox((0, 0), lbl, font=font_sm)
            tw2 = bbox2[2] - bbox2[0]
            draw.text((x + (stat_w - 30 - tw2) // 2, y + 140), lbl, fill=BRAND["text_muted"], font=font_sm)

    # Sections
    sections = content.get("infographic_spec", {}).get("sections", [])
    y_start = 560
    section_h = (H - y_start - 120) // max(len(sections), 1)
    for i, sec in enumerate(sections):
        y = y_start + i * section_h
        colour = colours[i % len(colours)]
        draw.rounded_rectangle([(60, y), (W - 60, y + section_h - 20)], radius=12, fill="white", outline=BRAND["border"])

        # Section label bar
        draw.rounded_rectangle([(60, y), (380, y + 50)], radius=8, fill=colour)
        draw.text((80, y + 8), sec.get("label", ""), fill="white", font=font_md)

        # Items
        items = sec.get("items", [])
        for j, item in enumerate(items[:6]):
            iy = y + 70 + j * 42
            draw.text((100, iy), f"•  {item[:90]}", fill=BRAND["text"], font=font_sm)

    # Footer
    draw.rectangle([(0, H - 80), (W, H)], fill=BRAND["primary"])
    draw.text((80, H - 60), "Amazon PPC Manager Training Program · Student Workbook · v2026", fill="#93c5fd", font=font_xs)

    img.save(str(out_path), "PNG", quality=95)


def generate_xlsx(content: dict, out_path: Path):
    """Generate an Excel workbook with calculators and data tables."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    primary_fill = PatternFill(start_color="1a56db", end_color="1a56db", fill_type="solid")
    accent_fill = PatternFill(start_color="f59e0b", end_color="f59e0b", fill_type="solid")
    light_fill = PatternFill(start_color="f0f4f8", end_color="f0f4f8", fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    title_font = Font(name="Calibri", size=14, bold=True, color="1a56db")
    border = Border(
        bottom=Side(style="thin", color="E5E7EB"),
        right=Side(style="thin", color="E5E7EB"),
    )

    sheets = content.get("spreadsheet_spec", {}).get("sheets", [])
    if not sheets:
        # Fallback sheet
        sheets = [{"name": "Calculator", "columns": ["Metric", "Value"], "sample_rows": [["ACoS", "25%"]], "formulas": {}, "notes": ""}]

    for idx, sheet_spec in enumerate(sheets):
        if idx == 0:
            ws = wb.active
            ws.title = sheet_spec.get("name", f"Sheet{idx+1}")
        else:
            ws = wb.create_sheet(sheet_spec.get("name", f"Sheet{idx+1}"))

        # Title
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=min(len(sheet_spec.get("columns", [])), 8))
        title_cell = ws.cell(row=1, column=1, value=f"{content.get('module_code', '')} — {sheet_spec.get('name', '')}")
        title_cell.font = title_font
        title_cell.alignment = Alignment(horizontal="left", vertical="center")
        ws.row_dimensions[1].height = 30

        # Headers
        columns = sheet_spec.get("columns", ["A", "B", "C"])
        for ci, col in enumerate(columns, 1):
            cell = ws.cell(row=3, column=ci, value=col)
            cell.font = header_font
            cell.fill = primary_fill
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell.border = border
        ws.row_dimensions[3].height = 28

        # Sample data
        rows = sheet_spec.get("sample_rows", [])
        for ri, row_data in enumerate(rows, 4):
            for ci, val in enumerate(row_data, 1):
                cell = ws.cell(row=ri, column=ci, value=val)
                cell.font = Font(name="Calibri", size=10)
                cell.alignment = Alignment(vertical="center", wrap_text=True)
                cell.border = border
                if ri % 2 == 0:
                    cell.fill = light_fill

        # Formulas
        formulas = sheet_spec.get("formulas", {})
        for cell_ref, formula in formulas.items():
            try:
                col_letter = re.match(r"([A-Z]+)", cell_ref).group(1)
                row_num = int(re.search(r"(\d+)$", cell_ref).group(1))
                ws[f"{cell_ref}"] = formula
            except (AttributeError, ValueError):
                pass

        # Auto-width
        for ci, col in enumerate(columns, 1):
            max_len = max(len(str(col)), max((len(str(r[ci-1])) if ci-1 < len(r) else 0) for r in rows) if rows else 0)
            ws.column_dimensions[get_column_letter(ci)].width = min(max_len + 4, 40)

        # Notes
        notes = sheet_spec.get("notes", "")
        if notes:
            note_row = len(rows) + 5
            ws.cell(row=note_row, column=1, value=f"Notes: {notes}").font = Font(name="Calibri", size=9, italic=True, color="6b7280")

    wb.save(str(out_path))


# ── Main ────────────────────────────────────────────────────────────────────
def list_topics():
    for phase_num, phase in TOPICS.items():
        print(f"\nPhase {phase_num}: {phase['title']}")
        for code, topic in phase["modules"]:
            print(f"  {code}  {topic}")


def main():
    parser = argparse.ArgumentParser(description="PPC Creative Crew — generate professional training materials")
    parser.add_argument("topic", nargs="?", help="Module topic")
    parser.add_argument("--phase", type=int, default=2, help="Phase (1-4)")
    parser.add_argument("--module", default="2.1", help="Module code")
    parser.add_argument("--api-key", help="OpenRouter API key")
    parser.add_argument("--model", default=LLM_MODEL)
    parser.add_argument("--output", "-o", default="./creative-output", help="Output directory")
    parser.add_argument("--list-topics", action="store_true")
    parser.add_argument("--skip-pptx", action="store_true")
    parser.add_argument("--skip-pdf", action="store_true")
    parser.add_argument("--skip-infographic", action="store_true")
    parser.add_argument("--skip-xlsx", action="store_true")
    args = parser.parse_args()

    if args.list_topics:
        list_topics()
        sys.exit(0)

    if not args.topic:
        print("ERROR: topic required. --list-topics for options.", file=sys.stderr)
        sys.exit(1)

    if args.api_key:
        os.environ["OPENAI_API_KEY"] = args.api_key
        os.environ["OPENROUTER_API_KEY"] = args.api_key

    llm = LLM(model=args.model, api_key=os.environ["OPENAI_API_KEY"], base_url=LLM_BASE)
    agents = make_agents(llm)
    phase = max(1, min(4, args.phase))
    code = args.module
    topic = args.topic

    tasks = make_tasks(agents, phase, code, topic)
    crew = Crew(agents=list(agents.values()), tasks=tasks, process=Process.sequential, verbose=False)

    print(f"Creative Crew: Phase {phase}, Module {code}: {topic}", file=sys.stderr)
    print(f"Model: {args.model}", file=sys.stderr)
    print(f"Team: Researcher, Writer, Infographic Designer, Spreadsheet Engineer, Creative Director", file=sys.stderr)
    print(file=sys.stderr)

    result = crew.kickoff()
    raw = str(result)
    content = extract_json(raw)
    if not content:
        print("ERROR: Could not parse crew output as JSON", file=sys.stderr)
        print(raw, file=sys.stderr)
        sys.exit(1)

    content["phase"] = phase
    content["module_code"] = code
    content["topic"] = topic

    out_dir = Path(args.output)
    out_dir.mkdir(parents=True, exist_ok=True)

    # Generate all outputs
    generated = []
    if not args.skip_pptx:
        pptx_path = out_dir / f"{code}-presentation.pptx"
        generate_pptx(content, pptx_path)
        generated.append(("PPTX", pptx_path))

    if not args.skip_pdf:
        pdf_path = out_dir / f"{code}-handout.pdf"
        actual = generate_pdf_handout(content, pdf_path)
        label = "PDF" if actual.suffix == ".pdf" else "HTML"
        generated.append((label, actual))

    if not args.skip_infographic:
        inf_path = out_dir / f"{code}-infographic.png"
        generate_infographic(content, inf_path)
        generated.append(("PNG", inf_path))

    if not args.skip_xlsx:
        xlsx_path = out_dir / f"{code}-workbook.xlsx"
        generate_xlsx(content, xlsx_path)
        generated.append(("XLSX", xlsx_path))

    # Save raw content JSON
    json_path = out_dir / f"{code}-content.json"
    json_path.write_text(json.dumps(content, indent=2))
    generated.append(("JSON", json_path))

    print(f"\nGenerated {len(generated)} files in {out_dir}/:", file=sys.stderr)
    for fmt, path in generated:
        size = path.stat().st_size
        print(f"  [{fmt}]  {path.name}  ({size:,} bytes)", file=sys.stderr)
    print(file=sys.stderr)
    print("Creative crew complete.", file=sys.stderr)


if __name__ == "__main__":
    main()
